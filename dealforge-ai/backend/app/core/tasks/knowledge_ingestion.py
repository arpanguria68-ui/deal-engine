"""
Knowledge Base Ingestion — Batch-ingest Excel, PDF, PPTX documents into PageIndex RAG.

Supports:
- Excel (.xlsx): Extracts sheet names, headers, cell-level formula descriptions
- PDF (.pdf): Text extraction via PyMuPDF
- PPTX (.pptx): Slide text extraction via python-pptx
"""

import os
import json
from typing import Dict, Any, Optional, List
from pathlib import Path
from datetime import datetime
import structlog

logger = structlog.get_logger()

# Default knowledge base paths
DEFAULT_KNOWLEDGE_DIRS = [
    r"F:\code project\Kimi_Agent_DealForge AI PRD\Knowledge managerment\Excel knowledge",
    r"F:\code project\Kimi_Agent_DealForge AI PRD\Knowledge managerment\Finance knowledge base",
]


class KnowledgeIngestionService:
    """Batch-ingest documents into the PageIndex RAG system."""

    def __init__(self, pageindex_client=None):
        self.pageindex = pageindex_client
        self.logger = structlog.get_logger()
        self._stats = {
            "total_files": 0,
            "ingested": 0,
            "failed": 0,
            "skipped": 0,
            "by_type": {"xlsx": 0, "pdf": 0, "pptx": 0, "other": 0},
        }

    # ── Excel Parsing ─────────────────────────────────────
    @staticmethod
    def parse_excel(file_path: str) -> str:
        """Extract structured content from an Excel file."""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            parts = [f"FILE: {os.path.basename(file_path)}"]
            parts.append(f"Sheets: {', '.join(wb.sheetnames)}\n")

            for sheet_name in wb.sheetnames[:10]:  # Limit to 10 sheets
                ws = wb[sheet_name]
                parts.append(f"\n=== SHEET: {sheet_name} ===")

                rows_seen = 0
                for row in ws.iter_rows(max_row=50, values_only=False):  # First 50 rows
                    cell_values = []
                    for cell in row[:20]:  # First 20 columns
                        if cell.value is not None:
                            cell_values.append(str(cell.value))
                    if cell_values:
                        parts.append(" | ".join(cell_values))
                        rows_seen += 1
                if rows_seen == 0:
                    parts.append("(empty sheet)")

            wb.close()
            return "\n".join(parts)[:10000]  # Cap at 10K chars

        except Exception as e:
            logger.warning("excel_parse_error", file=file_path, error=str(e))
            return f"FILE: {os.path.basename(file_path)}\nError parsing: {e}"

    # ── PDF Parsing ───────────────────────────────────────
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """Extract text from a PDF file."""
        try:
            import pymupdf

            doc = pymupdf.open(file_path)
            parts = [f"FILE: {os.path.basename(file_path)}"]
            parts.append(f"Pages: {len(doc)}\n")

            for i, page in enumerate(doc):
                if i >= 30:  # Cap at 30 pages
                    parts.append(f"\n... ({len(doc) - 30} more pages)")
                    break
                text = page.get_text()
                if text.strip():
                    parts.append(f"\n--- Page {i+1} ---\n{text[:2000]}")

            doc.close()
            return "\n".join(parts)[:15000]

        except Exception as e:
            logger.warning("pdf_parse_error", file=file_path, error=str(e))
            return f"FILE: {os.path.basename(file_path)}\nError: {e}"

    # ── PPTX Parsing ──────────────────────────────────────
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """Extract text from a PowerPoint file."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            parts = [f"FILE: {os.path.basename(file_path)}"]
            parts.append(f"Slides: {len(prs.slides)}\n")

            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    parts.append(f"\n--- Slide {i+1} ---")
                    parts.append("\n".join(texts[:20]))

            return "\n".join(parts)[:10000]

        except Exception as e:
            logger.warning("pptx_parse_error", file=file_path, error=str(e))
            return f"FILE: {os.path.basename(file_path)}\nError: {e}"

    # ── Batch Ingestion ───────────────────────────────────
    async def ingest_directory(
        self,
        directory: str,
        file_types: Optional[List[str]] = None,
    ) -> Dict[str, Any]:
        """Ingest all supported files from a directory into PageIndex."""
        if file_types is None:
            file_types = [".xlsx", ".xls", ".pdf", ".pptx", ".docx"]

        directory = Path(directory)
        if not directory.exists():
            return {"error": f"Directory not found: {directory}"}

        files = []
        for ext in file_types:
            files.extend(directory.rglob(f"*{ext}"))

        self._stats["total_files"] = len(files)
        self.logger.info(
            "ingestion_started", directory=str(directory), files=len(files)
        )

        for fpath in files:
            await self._ingest_file(str(fpath))

        self.logger.info("ingestion_completed", **self._stats)
        return {**self._stats, "directory": str(directory)}

    async def _ingest_file(self, file_path: str) -> bool:
        """Parse and ingest a single file."""
        ext = Path(file_path).suffix.lower()
        try:
            # Parse based on file type
            if ext in (".xlsx", ".xls"):
                content = self.parse_excel(file_path)
                self._stats["by_type"]["xlsx"] += 1
            elif ext == ".pdf":
                content = self.parse_pdf(file_path)
                self._stats["by_type"]["pdf"] += 1
            elif ext == ".pptx":
                content = self.parse_pptx(file_path)
                self._stats["by_type"]["pptx"] += 1
            else:
                self._stats["skipped"] += 1
                return False

            if not content or len(content) < 20:
                self._stats["skipped"] += 1
                return False

            # Ingest into PageIndex
            if self.pageindex:
                metadata = {
                    "source": os.path.basename(file_path),
                    "file_type": ext.lstrip("."),
                    "full_path": file_path,
                    "ingested_at": datetime.utcnow().isoformat(),
                }
                await self.pageindex.ingest_text(content, metadata=metadata)
                self._stats["ingested"] += 1
                return True
            else:
                # No pageindex client — just count
                self._stats["ingested"] += 1
                return True

        except Exception as e:
            self.logger.warning("file_ingestion_failed", file=file_path, error=str(e))
            self._stats["failed"] += 1
            return False

    async def ingest_all_knowledge_bases(self) -> Dict[str, Any]:
        """Ingest all default knowledge base directories."""
        all_results = {}
        for directory in DEFAULT_KNOWLEDGE_DIRS:
            if os.path.exists(directory):
                result = await self.ingest_directory(directory)
                all_results[directory] = result
            else:
                all_results[directory] = {"error": "Directory not found"}
        return all_results

    def get_stats(self) -> Dict[str, Any]:
        return {**self._stats}
