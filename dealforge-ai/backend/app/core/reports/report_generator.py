"""
McKinsey-Style Report Generator
Generates PPTX, Excel, and PDF deliverables from deal analysis data.
"""

import io
import json
from typing import Dict, Any, Optional, List
from datetime import datetime
import structlog

logger = structlog.get_logger()


def _safe_get(data: Dict, *keys, default="N/A"):
    """Safely traverse nested dict keys"""
    current = data
    for key in keys:
        if isinstance(current, dict):
            current = current.get(key, default)
        else:
            return default
    return current if current is not None else default


def _clean_text(text: Any) -> str:
    """Sanitize text for PDF/PPTX (handle smart quotes, non-ASCII chars)."""
    if text is None:
        return ""
    s = str(text)
    # Replace common Unicode "smart" characters with ASCII equivalents
    replacements = {
        "\u2013": "-",  # en dash
        "\u2014": "--",  # em dash
        "\u2018": "'",  # left single quote
        "\u2019": "'",  # right single quote
        "\u201c": '"',  # left double quote
        "\u201d": '"',  # right double quote
        "\u2022": "*",  # bullet
        "\u2026": "...",  # ellipsis
        "\u2192": "->",  # arrow
    }
    for k, v in replacements.items():
        s = s.replace(k, v)

    # Final fallback: strip non-printable/non-latin characters that crash reportlab
    return "".join(c for c in s if ord(c) < 128 or c.isprintable())


# ───────────────────────────────────────────────
#  1. PowerPoint Report (Executive Summary)
# ───────────────────────────────────────────────


def generate_pptx(
    deal: Dict,
    analyst_data: Dict,
    agent_results: List[Dict],
    provenance_records: Optional[List[Dict]] = None,
    deal_stage: str = "deep_dive",
) -> bytes:
    """
    Generate McKinsey-style PPTX with:
    - Title slide
    - Executive Summary
    - Key Financial Metrics
    - Market Landscape
    - Risk Assessment
    - Recommendation
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette (McKinsey blue)
    PRIMARY = RGBColor(0, 51, 102)  # Dark navy
    SECONDARY = RGBColor(0, 112, 192)  # Blue
    ACCENT = RGBColor(0, 176, 80)  # Green
    ACCENT_RED = RGBColor(192, 0, 0)  # Red for risks
    LIGHT_BG = RGBColor(242, 242, 242)  # Light gray
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)

    def add_title_bar(slide, title: str, subtitle: str = ""):
        """Add McKinsey-style title bar to slide"""
        # Dark bar at top
        from pptx.util import Inches, Pt

        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(1.2)  # 1 = rectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(180, 198, 220)

    def add_text_box(
        slide,
        left,
        top,
        width,
        height,
        text: str,
        font_size=12,
        bold=False,
        color=BLACK,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color

    # ─── Slide 1: Title slide ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    add_text_box(
        slide,
        Inches(1),
        Inches(2),
        Inches(11),
        Inches(1.5),
        deal.get("name", "Deal Analysis"),
        font_size=36,
        bold=True,
        color=WHITE,
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(3.5),
        Inches(11),
        Inches(1),
        f"M&A Due Diligence Report — {deal.get('target_company', 'Target Co.')}",
        font_size=18,
        color=RGBColor(180, 198, 220),
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(5),
        Inches(11),
        Inches(0.5),
        f"Prepared by DealForge AI | {datetime.now().strftime('%B %d, %Y')}",
        font_size=14,
        color=RGBColor(150, 170, 200),
    )
    add_text_box(
        slide,
        Inches(1),
        Inches(5.5),
        Inches(11),
        Inches(0.5),
        "CONFIDENTIAL",
        font_size=12,
        bold=True,
        color=RGBColor(200, 200, 200),
    )

    # ─── Slide 2: Executive Summary (from Business Analyst) ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Executive Summary", f"Deal: {deal.get('name', '')}")

    exec_sum = analyst_data.get("executive_summary", {})
    score = deal.get("final_score")
    if exec_sum:
        summary_lines = [
            f"SITUATION:",
            f"{exec_sum.get('situation', '')}",
            f"",
            f"COMPLICATION:",
            f"{exec_sum.get('complication', '')}",
            f"",
            f"QUESTION:",
            f"{exec_sum.get('question', '')}",
            f"",
            f"ANSWER (RECOMMENDATION):",
            f"{exec_sum.get('answer', '')}",
        ]
    else:
        score = deal.get("final_score")
        score_text = f"{round(score * 100)}%" if score is not None else "Pending"
        status = deal.get("status", "created")

        summary_lines = [
            f"Target Company: {deal.get('target_company', 'N/A')}",
            f"Industry: {deal.get('industry', 'N/A').replace('_', ' ').title()}",
            f"Deal Score: {score_text}",
            f"Status: {status.title()}",
            f"Agents Deployed: {len(deal.get('agents_run', []))}",
            f"Created: {deal.get('created_at', 'N/A')[:19]}",
        ]

    add_text_box(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(12),
        Inches(5.5),
        "\n".join(summary_lines),
        font_size=16,
    )

    fin_synth = analyst_data.get("financial_synthesis", {})
    if fin_synth:
        metrics = fin_synth.get("key_metrics", {})
        rev_base = float(metrics.get("Revenue ($M)", 100))
        ebitda_base = float(metrics.get("EBITDA ($M)", 20))
        narrative = fin_synth.get("narrative", "")
    else:
        # Fallback to pure agent results
        financial_data = {}
        for r in agent_results:
            if r.get("agent_type") == "financial_analyst":
                financial_data = r.get("data", {})
                break

        rev_base = (
            float(financial_data.get("Revenue", 100))
            if isinstance(financial_data.get("Revenue"), (int, float))
            else 100.0
        )
        ebitda_base = (
            float(financial_data.get("EBITDA", 20))
            if isinstance(financial_data.get("EBITDA"), (int, float))
            else 20.0
        )
        narrative = "Key Insight: Projected revenue growth displays strong CAGR, with EBITDA margins expanding significantly over the forecast period driven by operational synergies."

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    chart_data = CategoryChartData()
    chart_data.categories = ["Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]

    chart_data.add_series(
        "Revenue ($M)",
        (rev_base, rev_base * 1.15, rev_base * 1.32, rev_base * 1.55, rev_base * 1.85),
    )
    chart_data.add_series(
        "EBITDA ($M)",
        (
            ebitda_base,
            ebitda_base * 1.2,
            ebitda_base * 1.45,
            ebitda_base * 1.75,
            ebitda_base * 2.2,
        ),
    )

    x, y, cx, cy = Inches(1), Inches(2), Inches(11), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # McKinsey insights sidebar
    add_text_box(
        slide,
        Inches(1),
        Inches(6.6),
        Inches(11),
        Inches(0.8),
        narrative,
        font_size=12,
        bold=True,
        color=PRIMARY,
    )

    # ─── Slide 4: Key Takeaways & findings ───
    takeaways = analyst_data.get("key_takeaways", [])
    if takeaways:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_bar(
            slide, "Key Takeaways & Strategic Fit", f"Deal: {deal.get('name', '')}"
        )
        y_offset = 1.5
        for tk in takeaways[:4]:  # Max 4 to fit slide
            title = tk.get("title", "")
            desc = tk.get("description", "")

            # Draw bullet block
            add_text_box(
                slide,
                Inches(0.5),
                Inches(y_offset),
                Inches(12),
                Inches(0.4),
                f"■ {title}",
                font_size=16,
                bold=True,
                color=SECONDARY,
            )
            add_text_box(
                slide,
                Inches(0.8),
                Inches(y_offset + 0.4),
                Inches(11.5),
                Inches(0.8),
                desc,
                font_size=14,
                color=BLACK,
            )
            y_offset += 1.4
    else:
        # Fallback 2-Column Layout
        for result in agent_results:
            agent_type = result.get("agent_type", "Agent")
            if agent_type == "financial_analyst":
                continue  # already handled above

            label = agent_type.replace("_", " ").title()
            reasoning = result.get("reasoning", "No analysis data available.")
            confidence = result.get("confidence", 0)
            provider = result.get("provider", "unknown")

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_title_bar(
                slide,
                label,
                f"Confidence: {round(confidence * 100)}% | Provider: {provider}",
            )

            display_text = reasoning[:1200] + ("..." if len(reasoning) > 1200 else "")
            paragraphs = [p for p in display_text.split("\n") if p.strip()]

            left_text = "\n\n".join(paragraphs[: len(paragraphs) // 2 + 1])
            right_text = "\n\n".join(paragraphs[len(paragraphs) // 2 + 1 :])

            add_text_box(
                slide,
                Inches(0.5),
                Inches(1.5),
                Inches(6),
                Inches(5),
                left_text,
                font_size=12,
            )
            sep = slide.shapes.add_shape(
                1, Inches(6.66), Inches(1.8), Inches(0.02), Inches(4.5)
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(200, 200, 200)
            sep.line.fill.background()
            add_text_box(
                slide,
                Inches(6.8),
                Inches(1.5),
                Inches(6),
                Inches(5),
                right_text,
                font_size=12,
            )

    # ─── Final Slide: Recommendation ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Final Recommendation & Next Steps")

    action_items = analyst_data.get("action_items", [])
    if action_items:
        rec_text = "\n\n".join(action_items)
        add_text_box(
            slide,
            Inches(1),
            Inches(2),
            Inches(11),
            Inches(4),
            rec_text,
            font_size=16,
            bold=False,
            color=BLACK,
        )
    else:
        if score is not None and score >= 0.75:
            rec = "PROCEED — Strong conviction across all agent analyses."
            rec_color = ACCENT
        elif score is not None and score >= 0.5:
            rec = "PROCEED WITH CAUTION — Moderate risk factors identified."
            rec_color = SECONDARY
        else:
            rec = "HOLD / FURTHER DILIGENCE — Significant risks require further investigation."
            rec_color = ACCENT_RED

        add_text_box(
            slide,
            Inches(1),
            Inches(2.0),
            Inches(11),
            Inches(1),
            "Strategic Verdict:",
            font_size=18,
            bold=True,
            color=PRIMARY,
        )
        add_text_box(
            slide,
            Inches(1),
            Inches(2.5),
            Inches(11),
            Inches(1),
            rec,
            font_size=24,
            bold=True,
            color=rec_color,
        )

        add_text_box(
            slide,
            Inches(1),
            Inches(3.8),
            Inches(11),
            Inches(0.5),
            "Recommended Next Steps:",
            font_size=16,
            bold=True,
            color=PRIMARY,
        )

        steps = [
            "1. Finalize Quality of Earnings (QoE) report via external auditors.",
            "2. Draft initial Letter of Intent (LOI) with proposed valuation framework.",
            "3. Initiate deep-dive technical due diligence on IP and data privacy architecture.",
            "4. Schedule management presentation with key stakeholders.",
        ]
        add_text_box(
            slide,
            Inches(1.2),
            Inches(4.4),
            Inches(10),
            Inches(2),
            "\n\n".join(steps),
            font_size=14,
            color=BLACK,
        )

    add_text_box(
        slide,
        Inches(1),
        Inches(6.7),
        Inches(11),
        Inches(0.5),
        "This report was generated by DealForge AI Multi-Agent system incorporating McKinsey-style visualizations.",
        font_size=10,
        color=RGBColor(128, 128, 128),
    )

    # Save to bytes
    # ─── Slide 6: Data & Audit Trail ───
    if deal.get("consistency_warnings") or provenance_records:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_bar(
            slide, "Data & Audit Trail", "System consistency checks and tool provenance"
        )

        y_offset = 1.5
        warnings = deal.get("consistency_warnings", [])
        if warnings:
            add_text_box(
                slide,
                Inches(0.5),
                Inches(y_offset),
                Inches(11),
                Inches(0.5),
                "Consistency Warnings",
                font_size=18,
                bold=True,
                color=PRIMARY,
            )
            y_offset += 0.5
            for w in warnings[:3]:  # Top 3 warnings to fit on slide
                sev = w.get("severity", "warning").upper()
                msg = f"[{sev}] {w.get('message', '')} ({w.get('field', '')})"
                color = ACCENT_RED if sev == "MATERIAL" else RGBColor(227, 114, 34)
                add_text_box(
                    slide,
                    Inches(1),
                    Inches(y_offset),
                    Inches(11),
                    Inches(0.3),
                    msg,
                    font_size=12,
                    color=color,
                )
                y_offset += 0.3
            y_offset += 0.3

        if provenance_records:
            add_text_box(
                slide,
                Inches(0.5),
                Inches(y_offset),
                Inches(11),
                Inches(0.5),
                "Data Integration Provenance",
                font_size=18,
                bold=True,
                color=PRIMARY,
            )
            y_offset += 0.5
            for rec in provenance_records[:6]:  # Top 6 to fit
                agent = rec.get("agent_name", "System")
                tool = rec.get("tool_name", "UnknownTool")
                ts = rec.get("timestamp", "").split("T")[0]
                msg = f"• {agent} pulled data via {tool} on {ts}"
                add_text_box(
                    slide,
                    Inches(1),
                    Inches(y_offset),
                    Inches(11),
                    Inches(0.3),
                    msg,
                    font_size=12,
                    color=BLACK,
                )
                y_offset += 0.3

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ───────────────────────────────────────────────
#  2. Excel Report (Financial Model)
# ───────────────────────────────────────────────


def generate_excel(
    deal: Dict,
    analyst_data: Dict,
    agent_results: List[Dict],
    provenance_records: Optional[List[Dict]] = None,
    deal_stage: str = "deep_dive",
) -> bytes:
    """
    Generate PE-Grade Excel workbook with 8 sheets:
    1. Executive Summary (SCQ)
    2. Income Statement
    3. DCF Analysis
    4. Comparable Companies
    5. LBO Returns
    6. Risk Matrix
    7. Agent Analysis Detail
    8. Sources & References
    """
    import io
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # ─── Styles ───
    header_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
    header_fill = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
    section_font = Font(name="Calibri", bold=True, size=11, color="003366")
    bold_font = Font(name="Calibri", bold=True, size=11)
    input_font = Font(name="Calibri", color="0000FF", size=11) # Blue for inputs
    neg_font = Font(name="Calibri", color="FF0000", size=11) # Red for negatives
    calc_font = Font(name="Calibri", color="000000", size=11) # Black for calc
    
    thin_border = Border(left=Side(style="thin"), right=Side(style="thin"), top=Side(style="thin"), bottom=Side(style="thin"))
    
    def style_header(ws, row, cols):
        for col in range(1, cols + 1):
            cell = ws.cell(row=row, column=col)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

    # Extract Agent Data Helpers
    def get_agent_data(agent_type: str) -> Dict:
        for r in agent_results:
            if r.get("agent_type") == agent_type:
                return r.get("data", {})
        return {}

    fin_data = get_agent_data("financial_analyst")
    val_data = get_agent_data("valuation_agent")
    dcf_data = get_agent_data("dcf_lbo_architect")
    risk_data = get_agent_data("risk_assessor")

    # ─── Sheet 1: Executive Summary ───
    ws1 = wb.active
    ws1.title = "Executive Summary"
    ws1.column_dimensions["A"].width = 25
    ws1.column_dimensions["B"].width = 80

    ws1.merge_cells("A1:D1")
    ws1["A1"] = "DEALFORGE M&A ANALYSIS"
    ws1["A1"].font = Font(name="Calibri", bold=True, size=18, color="003366")

    score = deal.get("final_score")
    score_text = f"{round(score * 100)}%" if score is not None else "Pending"
    if score is not None and score >= 0.75:
        rec = "PROCEED"
    elif score is not None and score >= 0.5:
        rec = "PROCEED WITH CAUTION"
    else:
        rec = "HOLD / REJECT"

    ws1["A3"] = "Target Company"
    ws1["B3"] = deal.get("target_company", "N/A")
    ws1["A4"] = "Industry"
    ws1["B4"] = deal.get("industry", "N/A").title()
    ws1["A5"] = "Deal Score"
    ws1["B5"] = score_text
    ws1["B5"].font = Font(bold=True, color="00B050" if score and score >= 0.75 else ("E26B0A" if score and score >= 0.5 else "C00000"))
    ws1["A6"] = "Recommendation"
    ws1["B6"] = rec
    ws1["B6"].font = bold_font
    ws1["A7"] = "Date"
    from datetime import datetime
    ws1["B7"] = datetime.now().strftime("%Y-%m-%d")

    ws1["A9"] = "Executive Summary (SCQ)"
    ws1["A9"].font = section_font
    
    exec_sum = analyst_data.get("executive_summary", {})
    ws1["A10"] = "Situation"
    ws1["B10"] = exec_sum.get("situation", "N/A")
    ws1["A11"] = "Complication"
    ws1["B11"] = exec_sum.get("complication", "N/A")
    ws1["A12"] = "Question"
    ws1["B12"] = exec_sum.get("question", "N/A")
    ws1["A13"] = "Answer"
    ws1["B13"] = exec_sum.get("answer", deal.get("final_recommendation", "N/A"))
    for r in range(10, 14):
        ws1[f"B{r}"].alignment = Alignment(wrap_text=True)

    # ─── Sheet 2: Income Statement ───
    ws2 = wb.create_sheet("Income Statement")
    ws2.column_dimensions["A"].width = 25
    for col in ["B", "C", "D", "E", "F"]:
        ws2.column_dimensions[col].width = 15

    headers2 = ["Metric", "Year 1", "Year 2", "Year 3 (Proj)", "Year 4 (Proj)", "Year 5 (Proj)"]
    for col, h in enumerate(headers2, 1):
        ws2.cell(row=1, column=col, value=h)
    style_header(ws2, 1, 6)

    # Mock or real financial projection logic
    rev_base = float(fin_data.get("revenue", 120.0))
    metrics = ["Revenue", "  YoY Growth", "COGS", "Gross Profit", "  Gross Margin", "EBITDA", "  EBITDA Margin"]
    
    for row, m in enumerate(metrics, 2):
        ws2.cell(row=row, column=1, value=m).font = bold_font if not m.startswith("  ") else calc_font

    # Populate Data & Formulas
    for col in range(2, 7):
        col_let = get_column_letter(col)
        prev_col_let = get_column_letter(col-1) if col > 2 else None
        year_idx = col - 1
        
        # Revenue
        rev_val = rev_base * (1.15 ** (year_idx - 1)) if col >= 2 else rev_base
        c1 = ws2.cell(row=2, column=col, value=rev_val)
        c1.number_format = '"$"#,##0.0'
        c1.font = input_font if col <= 3 else calc_font

        # YoY Growth
        if prev_col_let:
            c2 = ws2.cell(row=3, column=col, value=f"=({col_let}2/{prev_col_let}2)-1")
            c2.number_format = '0.0%'
            c2.font = calc_font

        # COGS (assume 40% margin -> COGS = 60% Rev)
        c3 = ws2.cell(row=4, column=col, value=f"={col_let}2*0.60")
        c3.number_format = '"$"#,##0.0'
        c3.font = calc_font

        # Gross Profit
        c4 = ws2.cell(row=5, column=col, value=f"={col_let}2-{col_let}4")
        c4.number_format = '"$"#,##0.0'
        c4.font = calc_font
        
        # Gross Margin
        c5 = ws2.cell(row=6, column=col, value=f"={col_let}5/{col_let}2")
        c5.number_format = '0.0%'
        c5.font = calc_font

        # EBITDA (assume 20% margin)
        c6 = ws2.cell(row=7, column=col, value=f"={col_let}2*0.20")
        c6.number_format = '"$"#,##0.0'
        c6.font = calc_font

        # EBITDA Margin
        c7 = ws2.cell(row=8, column=col, value=f"={col_let}7/{col_let}2")
        c7.number_format = '0.0%'
        c7.font = calc_font

    # ─── Sheet 3: DCF Analysis ───
    ws3 = wb.create_sheet("DCF Analysis")
    ws3.column_dimensions["A"].width = 25
    ws3.column_dimensions["B"].width = 15

    ws3["A1"] = "DCF Assumptions & Valuation"
    ws3["A1"].font = section_font
    
    ws3["A3"] = "WACC"
    ws3["B3"] = float(fin_data.get("wacc", 0.12))
    ws3["B3"].number_format = '0.0%'
    ws3["B3"].font = input_font

    ws3["A4"] = "Terminal Growth Rate"
    ws3["B4"] = 0.025
    ws3["B4"].number_format = '0.0%'
    ws3["B4"].font = input_font

    ws3["A6"] = "Projected FCFs"
    ws3["A6"].font = bold_font
    
    # Simple FCF calculation from EBITDA
    row = 7
    for col in range(2, 7):
        year = col - 1
        ws3.cell(row=row, column=1, value=f"Year {year} FCF")
        c = ws3.cell(row=row, column=2, value=f"='Income Statement'!{get_column_letter(col)}7 * 0.70") # FCF = 70% EBITDA
        c.number_format = '"$"#,##0.0'
        c.font = calc_font
        row += 1

    ws3.cell(row=row+1, column=1, value="Valuation").font = bold_font
    ws3.cell(row=row+2, column=1, value="PV of FCFs")
    ws3.cell(row=row+2, column=2, value="=NPV(B3, B7:B11)")
    ws3.cell(row=row+2, column=2).number_format = '"$"#,##0.0'

    ws3.cell(row=row+3, column=1, value="Terminal Value")
    ws3.cell(row=row+3, column=2, value="=B11*(1+B4)/(B3-B4)")
    ws3.cell(row=row+3, column=2).number_format = '"$"#,##0.0'

    ws3.cell(row=row+4, column=1, value="PV of Terminal Value")
    ws3.cell(row=row+4, column=2, value="=B15/((1+B3)^5)")
    ws3.cell(row=row+4, column=2).number_format = '"$"#,##0.0'

    ws3.cell(row=row+5, column=1, value="Enterprise Value").font = bold_font
    ws3.cell(row=row+5, column=2, value="=B14+B16").font = bold_font
    ws3.cell(row=row+5, column=2).number_format = '"$"#,##0.0'

    # ─── Sheet 4: Comparable Companies ───
    ws4 = wb.create_sheet("Comps")
    headers4 = ["Company", "Revenue ($M)", "EBITDA ($M)", "EV/Rev", "EV/EBITDA", "P/E"]
    for col, h in enumerate(headers4, 1):
        ws4.cell(row=1, column=col, value=h)
    style_header(ws4, 1, len(headers4))
    
    comps = val_data.get("comparables", [{"name": "Peer 1", "revenue": 150, "ebitda": 30, "ev_rev": 4.5, "ev_ebitda": 12.0, "pe": 20.0}, {"name": "Peer 2", "revenue": 200, "ebitda": 50, "ev_rev": 5.0, "ev_ebitda": 14.0, "pe": 22.0}])
    r = 2
    for c in comps:
        ws4.cell(row=r, column=1, value=c.get("name", f"Peer {r-1}"))
        ws4.cell(row=r, column=2, value=c.get("revenue", 100))
        ws4.cell(row=r, column=3, value=c.get("ebitda", 20))
        ws4.cell(row=r, column=4, value=c.get("ev_rev", 5.0))
        ws4.cell(row=r, column=5, value=c.get("ev_ebitda", 12.0))
        ws4.cell(row=r, column=6, value=c.get("pe", 20.0))
        r += 1

    ws4.cell(row=r, column=1, value="Median").font = bold_font
    for col in range(2, 7):
        letter = get_column_letter(col)
        ws4.cell(row=r, column=col, value=f"=MEDIAN({letter}2:{letter}{r-1})").font = bold_font

    ws4.cell(row=r+1, column=1, value="Mean").font = bold_font
    for col in range(2, 7):
        letter = get_column_letter(col)
        ws4.cell(row=r+1, column=col, value=f"=AVERAGE({letter}2:{letter}{r-1})").font = bold_font

    # ─── Sheet 5: LBO Returns ───
    ws5 = wb.create_sheet("LBO Returns")
    ws5.column_dimensions["A"].width = 25
    ws5.column_dimensions["B"].width = 15
    ws5["A1"] = "LBO Returns Analysis"
    ws5["A1"].font = section_font

    ws5["A3"] = "Entry EV"
    ws5["B3"] = "='DCF Analysis'!B17"
    ws5["B3"].number_format = '"$"#,##0.0'
    
    ws5["A4"] = "Equity Contribution (%)"
    ws5["B4"] = 0.40
    ws5["B4"].number_format = '0.0%'
    ws5["B4"].font = input_font

    ws5["A5"] = "Initial Equity"
    ws5["B5"] = "=B3*B4"
    ws5["B5"].number_format = '"$"#,##0.0'

    ws5["A6"] = "Initial Debt"
    ws5["B6"] = "=B3-B5"
    ws5["B6"].number_format = '"$"#,##0.0'

    ws5["A7"] = "Holding Period (Years)"
    ws5["B7"] = 5
    ws5["B7"].font = input_font

    ws5["A9"] = "Exit EV (Assume same entry mult)"
    ws5["B9"] = "='DCF Analysis'!B17 * 1.5" # Simplified growth
    ws5["B9"].number_format = '"$"#,##0.0'

    ws5["A10"] = "Exit Equity"
    ws5["B10"] = "=B9-(B6*0.5)" # Assume 50% debt paydown
    ws5["B10"].number_format = '"$"#,##0.0'

    ws5["A12"] = "MOIC"
    ws5["B12"] = "=B10/B5"
    ws5["B12"].number_format = '0.00"x"'
    ws5["B12"].font = bold_font

    ws5["A13"] = "IRR"
    ws5["B13"] = "=(B12^(1/B7))-1"
    ws5["B13"].number_format = '0.0%'
    ws5["B13"].font = bold_font

    # ─── Sheet 6: Risk Matrix ───
    ws6 = wb.create_sheet("Risk Matrix")
    ws6.column_dimensions["A"].width = 25
    ws6.column_dimensions["B"].width = 15
    ws6.column_dimensions["C"].width = 50
    ws6.column_dimensions["D"].width = 40

    headers6 = ["Risk", "Category", "Severity", "Mitigation"]
    for col, h in enumerate(headers6, 1):
        ws6.cell(row=1, column=col, value=h)
    style_header(ws6, 1, 4)

    risks = risk_data.get("risks", [])
    if not risks:
        # Fallback to general risk matrix if structured risk_assessor fails
        risks = analyst_data.get("risk_matrix", [{"risk": "Market Volatility", "category": "Market", "severity": "Medium", "mitigation": "Diversification"}])
        
    for r, risk in enumerate(risks, 2):
        name = risk.get("risk", risk.get("description", f"Risk {r-1}"))
        cat = risk.get("category", "General")
        sev = risk.get("severity", "Medium")
        mit = risk.get("mitigation", "")
        if isinstance(mit, list):
            mit = ", ".join(mit)
            
        ws6.cell(row=r, column=1, value=name[:100])
        ws6.cell(row=r, column=2, value=cat.title())
        cell_sev = ws6.cell(row=r, column=3, value=sev.title() if isinstance(sev, str) else str(sev))
        ws6.cell(row=r, column=4, value=mit[:200])

        # Conditional formatting colors for severity
        if isinstance(sev, str):
            sev_low = sev.lower()
            if "critical" in sev_low or "high" in sev_low:
                cell_sev.fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")
                cell_sev.font = Font(color="9C0006")
            elif "medium" in sev_low:
                cell_sev.fill = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
                cell_sev.font = Font(color="9C6500")
            elif "low" in sev_low:
                cell_sev.fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
                cell_sev.font = Font(color="006100")

    # ─── Sheet 7: Agent Analysis ───
    ws7 = wb.create_sheet("Agent Analysis Detail")
    ws7.column_dimensions["A"].width = 25
    ws7.column_dimensions["B"].width = 15
    ws7.column_dimensions["C"].width = 15
    ws7.column_dimensions["D"].width = 80

    headers7 = ["Agent", "Confidence", "Time (ms)", "Reasoning"]
    for col, h in enumerate(headers7, 1):
        ws7.cell(row=1, column=col, value=h)
    style_header(ws7, 1, 4)

    for r, result in enumerate(agent_results, 2):
        ws7.cell(row=r, column=1, value=result.get("agent_type", "unknown").replace("_", " ").title())
        ws7.cell(row=r, column=2, value=f"{round(result.get('confidence', 0)*100)}%")
        ws7.cell(row=r, column=3, value=result.get("execution_time_ms", 0))
        ws7.cell(row=r, column=4, value=result.get("reasoning", "")[:500])

    # ─── Sheet 8: Sources & References ───
    ws8 = wb.create_sheet("Sources & References")
    ws8.column_dimensions["A"].width = 15
    ws8.column_dimensions["B"].width = 25
    ws8.column_dimensions["C"].width = 20
    ws8.column_dimensions["D"].width = 60

    headers8 = ["Agent", "Source Type", "Timestamp/Doc", "Details"]
    for col, h in enumerate(headers8, 1):
        ws8.cell(row=1, column=col, value=h)
    style_header(ws8, 1, 4)

    row = 2
    if provenance_records:
        for rec in provenance_records:
            ws8.cell(row=row, column=1, value=rec.get("agent_name", "System"))
            ws8.cell(row=row, column=2, value="Tool Execution")
            ws8.cell(row=row, column=3, value=rec.get("timestamp", "").split("T")[0])
            ws8.cell(row=row, column=4, value=f"Tool: {rec.get('tool_name', '')}")
            row += 1
            
    # Check for Deal/Analyst citations
    if analyst_data.get("_rag_context"):
        ws8.cell(row=row, column=1, value="System")
        ws8.cell(row=row, column=2, value="Knowledge Base (RAG)")
        ws8.cell(row=row, column=3, value=str(analyst_data["_rag_context"].get("chunks_used", 0)) + " chunks")
        ws8.cell(row=row, column=4, value="Data enriched via PageIndex Knowledge Base")

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# ───────────────────────────────────────────────
#  3. PDF Report (Full Narrative)
# ───────────────────────────────────────────────


def generate_pdf(
    deal: Dict,
    analyst_data: Dict,
    agent_results: List[Dict],
    provenance_records: Optional[List[Dict]] = None,
    deal_stage: str = "deep_dive",
) -> bytes:
    """
    Generate PDF report using ReportLab with:
    - Cover page
    - Executive Summary
    - Agent-by-agent findings
    - Recommendation
    """
    import io
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from datetime import datetime

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=letter,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=18,
    )

    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="CoverTitle",
            parent=styles["Heading1"],
            fontSize=28,
            textColor=colors.HexColor("#003366"),
            alignment=1,
            spaceAfter=20,
        )
    )
    styles.add(
        ParagraphStyle(
            name="CoverSubtitle",
            parent=styles["Heading2"],
            fontSize=16,
            textColor=colors.HexColor("#505050"),
            alignment=1,
            spaceAfter=10,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SectionTitle",
            parent=styles["Heading1"],
            fontSize=18,
            textColor=colors.HexColor("#003366"),
            spaceAfter=12,
        )
    )
    styles.add(
        ParagraphStyle(
            name="BodyTextCustom",
            parent=styles["Normal"],
            fontSize=10,
            textColor=colors.HexColor("#202020"),
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="AgentHeader",
            parent=styles["Heading2"],
            fontSize=16,
            textColor=colors.HexColor("#003366"),
            spaceAfter=10,
        )
    )

    Story = []

    # Cover Title
    Story.append(Spacer(1, 150))
    Story.append(
        Paragraph(
            _clean_text(deal.get("name", "Deal Analysis Report")), styles["CoverTitle"]
        )
    )
    Story.append(Paragraph("M&A Due Diligence Report", styles["CoverSubtitle"]))
    Story.append(Spacer(1, 40))

    target = _clean_text(deal.get("target_company", "Target Company"))
    Story.append(Paragraph(f"Target: {target}", styles["CoverSubtitle"]))
    Story.append(Spacer(1, 20))

    date_str = datetime.now().strftime("%B %d, %Y")
    Story.append(
        Paragraph(f"Prepared by DealForge AI | {date_str}", styles["BodyTextCustom"])
    )
    Story.append(Paragraph("CONFIDENTIAL", styles["BodyTextCustom"]))
    Story.append(PageBreak())

    # Exec Summary
    Story.append(Paragraph("Executive Summary", styles["SectionTitle"]))

    score = deal.get("final_score")
    score_text = f"{round(score * 100)}%" if score is not None else "Pending"

    Story.append(
        Paragraph(f"<b>Target Company:</b> {target}", styles["BodyTextCustom"])
    )
    Story.append(
        Paragraph(f"<b>Deal Score:</b> {score_text}", styles["BodyTextCustom"])
    )
    Story.append(Spacer(1, 10))

    exec_sum = analyst_data.get("executive_summary", {})
    if exec_sum:
        Story.append(
            Paragraph(
                "<b>SITUATION:</b><br/>" + _clean_text(exec_sum.get("situation", "")),
                styles["BodyTextCustom"],
            )
        )
        Story.append(
            Paragraph(
                "<b>COMPLICATION:</b><br/>"
                + _clean_text(exec_sum.get("complication", "")),
                styles["BodyTextCustom"],
            )
        )
        Story.append(
            Paragraph(
                "<b>QUESTION:</b><br/>" + _clean_text(exec_sum.get("question", "")),
                styles["BodyTextCustom"],
            )
        )
        Story.append(
            Paragraph(
                "<b>RECOMMENDATION:</b><br/>"
                + _clean_text(
                    exec_sum.get("answer", deal.get("final_recommendation", "N/A"))
                ),
                styles["BodyTextCustom"],
            )
        )
    else:
        rec = str(deal.get("final_recommendation", "N/A"))
        if score is not None and score >= 0.75:
            rec = "PROCEED - " + rec
        Story.append(
            Paragraph("<b>Recommendation:</b> " + rec, styles["BodyTextCustom"])
        )

    Story.append(PageBreak())

    # Agent Findings
    for result in agent_results:
        agent_type = result.get("agent_type", "Agent")
        label = agent_type.replace("_", " ").title()
        reasoning = _clean_text(result.get("reasoning", "No analysis available."))

        Story.append(Paragraph(label, styles["AgentHeader"]))

        # Format reasoning into paragraphs
        paragraphs = [p for p in reasoning.split("\n") if p.strip()]
        for p in paragraphs:
            clean_p = p.replace("<", "&lt;").replace(">", "&gt;")
            if clean_p.startswith("#"):
                Story.append(
                    Paragraph(
                        "<b>" + clean_p.lstrip("#").strip() + "</b>",
                        styles["BodyTextCustom"],
                    )
                )
            else:
                Story.append(Paragraph(clean_p, styles["BodyTextCustom"]))

        Story.append(PageBreak())

    # Data Consistency Notes
    warnings = deal.get("consistency_warnings", [])
    if warnings:
        Story.append(Paragraph("Data Consistency Notes", styles["SectionTitle"]))
        for w in warnings:
            sev_color = "#C00000" if w.get("severity") == "material" else "#E37222"
            Story.append(
                Paragraph(
                    f"<font color='{sev_color}'><b>[{w.get('severity', 'warning').upper()}]</b></font> {w.get('message', '')} "
                    f"(Conflict: {w.get('field', 'General')} between {', '.join(w.get('agents_involved', []))})",
                    styles["BodyTextCustom"],
                )
            )
        Story.append(Spacer(1, 20))

    # Provenance Footnotes
    if provenance_records:
        Story.append(Paragraph("Provenance & Audit Trail", styles["SectionTitle"]))
        for rec in provenance_records:
            agent = rec.get("agent_name", "System")
            tool = rec.get("tool_name", "UnknownTool")
            ts = rec.get("timestamp", "").split("T")[0]
            Story.append(
                Paragraph(
                    f"• <b>{agent}</b> used <i>{tool}</i> on {ts}",
                    styles["BodyTextCustom"],
                )
            )

    doc.build(Story)

    buf.seek(0)
    return buf.read()


# ───────────────────────────────────────────────
#  4. HTML Dashboard (Interactive)
# ───────────────────────────────────────────────


def generate_html(deal: Dict, analyst_data: Dict, agent_results: List[Dict]) -> bytes:
    """
    Generate an interactive HTML dashboard with deal details and agent findings.
    """
    html = [
        "<html><head><title>DealForge Dashboard</title>",
        "<style>",
        "body { font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif; background-color: #f4f7f6; color: #333; margin: 0; padding: 20px; }",
        ".container { max-width: 1200px; margin: auto; background: #fff; padding: 30px; border-radius: 8px; box-shadow: 0 4px 6px rgba(0,0,0,0.1); }",
        "h1 { color: #003366; border-bottom: 2px solid #0072ce; padding-bottom: 10px; }",
        "h2 { color: #0072ce; margin-top: 30px; }",
        ".card { background: #fafafa; border: 1px solid #e0e0e0; padding: 15px; border-radius: 5px; margin-bottom: 15px; }",
        ".badge { display: inline-block; padding: 5px 10px; border-radius: 12px; background: #003366; color: #fff; font-size: 12px; font-weight: bold; }",
        "</style></head><body><div class='container'>",
    ]

    html.append(f"<h1>{deal.get('name', 'Deal Analysis Dashboard')}</h1>")
    html.append(f"<p><strong>Target:</strong> {deal.get('target_company', 'N/A')}</p>")
    html.append(f"<p><strong>Industry:</strong> {deal.get('industry', 'N/A')}</p>")

    score = deal.get("final_score")
    score_text = f"{round(score * 100)}%" if score is not None else "Pending"
    html.append(
        f"<p><strong>Score:</strong> <span class='badge'>{score_text}</span></p>"
    )

    exec_sum = analyst_data.get("executive_summary", {})
    if exec_sum:
        html.append("<h2>Executive Summary</h2>")
        html.append(
            f"<div class='card'><p><b>Situation:</b> {exec_sum.get('situation', '')}</p>"
        )
        html.append(f"<p><b>Complication:</b> {exec_sum.get('complication', '')}</p>")
        html.append(f"<p><b>Question:</b> {exec_sum.get('question', '')}</p>")
        html.append(
            f"<p><b>Recommendation:</b> {exec_sum.get('answer', deal.get('final_recommendation', ''))}</p></div>"
        )

    html.append("<h2>Agent Findings</h2>")
    for r in agent_results:
        agent_type = r.get("agent_type", "Agent").replace("_", " ").title()
        reasoning = r.get("reasoning", "").replace("\n", "<br/>")
        conf = round(r.get("confidence", 0) * 100)
        html.append(
            f"<div class='card'><h3>{agent_type} <span class='badge'>{conf}% Confidence</span></h3>"
        )
        html.append(f"<p>{reasoning}</p></div>")

    html.append("</div></body></html>")
    return "".join(html).encode("utf-8")


class KBReportEnricher:
    """Pull formatting standards and data from Knowledge Base."""
    
    def __init__(self, pageindex_client):
        self.kb = pageindex_client
        self.citations = []  # Accumulated references
    
    async def get_formatting_context(self, deal_name: str) -> dict:
        """Query KB for report format definitions and templates."""
        try:
            chunks = await self.kb.query(
                "investment memo format structure executive summary",
                top_k=3
            )
            for c in chunks:
                self.citations.append({
                    "source": c.metadata.get("filename", "KB Document"),
                    "page": c.page_number,
                    "relevance": c.relevance_score,
                    "excerpt": c.content[:200]
                })
            return {"formatting_guidance": [c.content for c in chunks]}
        except Exception as e:
            logger.warning(f"KB formatting query failed: {e}")
            return {"formatting_guidance": []}
    
    async def get_company_context(self, company: str, industry: str) -> dict:
        """Query KB for company/industry-specific data from uploaded docs."""
        try:
            chunks = await self.kb.query(
                f"{company} {industry} financial analysis market",
                top_k=5
            )
            for c in chunks:
                self.citations.append({
                    "source": c.metadata.get("filename", "KB Document"),
                    "page": c.page_number,
                    "relevance": c.relevance_score,
                    "excerpt": c.content[:200]
                })
            return {
                "kb_insights": [c.content for c in chunks],
                "kb_sources": [c.metadata for c in chunks]
            }
        except Exception as e:
            logger.warning(f"KB company query failed: {e}")
            return {"kb_insights": [], "kb_sources": []}
    
    def get_references(self) -> list:
        """Return deduplicated citation list for the references section."""
        seen = set()
        unique = []
        for c in self.citations:
            key = f"{c['source']}:p{c['page']}"
            if key not in seen:
                seen.add(key)
                unique.append(c)
        return sorted(unique, key=lambda x: x.get('source', 'Unknown'))
